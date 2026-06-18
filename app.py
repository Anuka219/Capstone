'''
FastAPI backend for the MEM & MIM Guide Bot.

What this server does:
1. Serves MEM-Guide-Bot-Website.html at http://localhost:8000/
2. Accepts chat messages at POST /api/chat
3. Sends the message to OpenAI with a strict HS Pforzheim grounding prompt
4. Sends the generated answer to ElevenLabs text-to-speech
5. Returns both text and base64 audio to the browser

Setup:
    python3 -m venv .venv
    source .venv/bin/activate
    pip install fastapi uvicorn httpx python-dotenv

Environment variables:
    OPENROUTER_API_KEY=your_openrouter_key
    ELEVENLABS_API_KEY=your_elevenlabs_key

Optional environment variables:
    OPENROUTER_MODEL=openai/gpt-4o
    ELEVENLABS_VOICE_ID=21m00Tcm4TlvDq8ikWAM
    ELEVENLABS_MODEL_ID=eleven_multilingual_v2
    PORT=8000

Run:
    uvicorn app:app --host 127.0.0.1 --port 8000 --reload

Important:
    The default ElevenLabs voice ID, 21m00Tcm4TlvDq8ikWAM, is a temporary
    pre-made placeholder voice. Replace ELEVENLABS_VOICE_ID with the cloned
    professor voice ID later.
'''

from __future__ import annotations

import base64
import asyncio
import hashlib
import json
import os
import re
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import List, Optional
from xml.etree import ElementTree

import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from pydantic import BaseModel, Field


load_dotenv()

ROOT = Path(__file__).resolve().parent
HTML_FILE = ROOT / "index.html"
MEMORY_DIR = ROOT / "bot_memory"
LESSONS_FILE = MEMORY_DIR / "learned_corrections.json"
CONVERSATION_FILE = MEMORY_DIR / "conversation_history.json"
# Short-term chat memory so the bot can follow up across turns. Kept small to
# stay well under Groq's per-minute token limit.
MAX_STORED_TURNS = 8           # messages kept per session (≈4 exchanges)
HISTORY_CONTEXT_MESSAGES = 2   # last exchange only — enough for follow-ups, fewer tokens
HISTORY_MSG_CHAR_CAP = 600     # truncate long past answers to save tokens
KNOWLEDGE_DIR = ROOT / "knowledge_docs"
TIMETABLE_FILE = KNOWLEDGE_DIR / "Semesteruebersicht_SS26.md"
TIMETABLE_FILE_PATTERN = "*Timetable*_SS26.md"
COURSE_PROFESSORS_FILE = KNOWLEDGE_DIR / "Course_Professors_SS26.md"

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_MODEL = os.getenv("OPENROUTER_MODEL", os.getenv("OPENAI_MODEL", "openai/gpt-4o"))
# The OpenRouter fallback key has no credits, so it only wastes time and leaks a
# raw "purchase credits" error to the user. Keep it OFF unless explicitly enabled
# with a funded key (set OPENROUTER_ENABLED=true in .env).
OPENROUTER_ENABLED = (os.getenv("OPENROUTER_ENABLED", "false").lower() in {"1", "true", "yes", "on"})

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
# Google Gemini as a reliable fallback brain. Its free tier has a far higher
# per-minute token limit than Groq, so it covers Groq's rate-limit overflow.
# Uses Gemini's OpenAI-compatible endpoint, so it slots in like the other providers.
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.0-flash")

ELEVENLABS_API_KEY = os.getenv("ELEVENLABS_API_KEY", "")
ELEVENLABS_VOICE_ID = os.getenv("ELEVENLABS_VOICE_ID") or "21m00Tcm4TlvDq8ikWAM"
ELEVENLABS_MODEL_ID = os.getenv("ELEVENLABS_MODEL_ID") or "eleven_turbo_v2_5"
# Lock the spoken language so the cloned (German) voice can't auto-switch to
# German mid-answer. Only turbo_v2_5 / flash_v2_5 / v3 honor language_code;
# multilingual_v2 ignores it and just auto-detects (the old buggy behavior).
ELEVENLABS_LANGUAGE = (os.getenv("ELEVENLABS_LANGUAGE") or "en").strip()
ELEVENLABS_LANGUAGE_LOCK_MODELS = {
    "eleven_turbo_v2_5",
    "eleven_flash_v2_5",
    "eleven_v3",
}
ELEVENLABS_OUTPUT_FORMAT = os.getenv("ELEVENLABS_OUTPUT_FORMAT") or "mp3_44100_128"
AUDIO_MIME = "audio/mpeg"
PROFESSOR_DEMO_AUDIO_FILE = os.getenv("PROFESSOR_DEMO_AUDIO_FILE", "")
VOICE_PROVIDER = (os.getenv("VOICE_PROVIDER") or "browser").lower()
XTTS_SPEAKER_WAV = os.getenv(
    "XTTS_SPEAKER_WAV",
    "voice_source/instant_voice_clone_clips/raphael_10s_45min.mp3",
)
XTTS_LANGUAGE = os.getenv("XTTS_LANGUAGE", "en")
XTTS_MAX_CHARS = int(os.getenv("XTTS_MAX_CHARS", "650"))
# Spoken-text cap for fast hosted providers (ElevenLabs/Cartesia/Fish). Much
# higher than the local-XTTS cap so the whole answer is read, not cut off.
HOSTED_MAX_CHARS = int(os.getenv("HOSTED_MAX_CHARS", "1500"))
XTTS_OUTPUT_DIR = ROOT / "voice_source/generated/chat_responses"

# HuggingFace OpenVoice V2 Space (primary voice when VOICE_PROVIDER=openvoice_hf).
# This clones the professor's voice on HuggingFace's servers — fast and good
# quality, but needs internet. Falls back to local XTTS automatically if down.
OPENVOICE_HF_SPACE = os.getenv("OPENVOICE_HF_SPACE", "myshell-ai/OpenVoiceV2")
OPENVOICE_STYLE = os.getenv("OPENVOICE_STYLE", "en_default")
OPENVOICE_FN_INDEX = int(os.getenv("OPENVOICE_FN_INDEX", "1"))
OPENVOICE_OUTPUT_DIR = ROOT / "voice_source/generated/chat_responses"

# Cartesia hosted voice cloning (VOICE_PROVIDER=cartesia). Fast (~1s, runs on
# their GPUs) and free-tier friendly. The bot auto-clones the professor clip on
# first use and caches the resulting voice id. Falls back to local XTTS if the
# key is missing or the API fails.
CARTESIA_API_KEY = os.getenv("CARTESIA_API_KEY", "")
CARTESIA_VOICE_ID = os.getenv("CARTESIA_VOICE_ID", "")
CARTESIA_MODEL = os.getenv("CARTESIA_MODEL", "sonic-2")
CARTESIA_VERSION = os.getenv("CARTESIA_VERSION", "2024-11-13")
CARTESIA_LANGUAGE = os.getenv("CARTESIA_LANGUAGE", "en")
CARTESIA_VOICE_ID_FILE = ROOT / "voice_source/cartesia_voice_id.txt"

# Fish Audio hosted cloning (VOICE_PROVIDER=fish). Fast, but the API is billed
# (needs API credit at fish.audio/app/developers). Auto-clones the professor
# clip into a Fish "model" on first use and caches the model id.
FISH_API_KEY = os.getenv("FISH_API_KEY", "")
FISH_REFERENCE_ID = os.getenv("FISH_REFERENCE_ID", "")
FISH_FORMAT = os.getenv("FISH_FORMAT", "mp3")
FISH_REFERENCE_ID_FILE = ROOT / "voice_source/fish_reference_id.txt"

MAX_LESSONS = 30
learned_corrections: List[dict] = []
knowledge_chunks: List[dict] = []
xtts_model = None
openvoice_hf_client = None
# XTTS shares one model object and is NOT thread/concurrency safe. Serialize all
# voice synthesis so two clicks can't run at once (which fails and then falls
# back to the static sample, making every answer sound identical).
voice_lock = asyncio.Lock()


SYSTEM_PROMPT = """
You are the MEM & MIM Guide Bot for Hochschule Pforzheim — a direct, knowledgeable study advisor for two master's programs: MEM (Engineering and Management, M.Sc.) and MIM (Industrial Management, M.Sc.).

FACTS — MEM (Engineering and Management, M.Sc.):
- Degree: M.Sc. Engineering and Management. Accredited by Akkreditierungsrat (national) and AACSB (international).
- 3 semesters incl. thesis (optional double degree: 4). 90 ECTS. 24 places. Starts summer and winter semester.
- Language of instruction: English and German.
- Deadlines: 15 January (summer), 15 June (winter, e.g. WS 2026/27).
- Admission: a degree graded "good" (2.5 or better) in industrial engineering / Wirtschaftsingenieurwesen; motivation letter (max 2 pages); a recommendation letter; English B2; and German C1 (e.g. DSH-2, TestDaF TDN 4x4, telc C1 Hochschule) ONLY for applicants who are not native German speakers and whose prior degree was not taught in German. C2 is NOT required; native German speakers (or a German-taught prior degree) submit no German certificate.
- Courses: Sem 1 — Leadership, Production Strategy, Process Management, + one elective (Engineering/Business/Design). Sem 2 — Market-oriented Product Development, Strategic Procurement, Value-oriented Management, Management of Emerging Technologies, Interdisciplinary Research Project. Sem 3 — Capstone Seminar, Scientific Colloquium, thesis.
- Careers: production/operations manager, process/industrial engineer moving into management, product manager, project manager (manufacturing), strategic procurement/supply-chain manager, technology & innovation manager, R&D/engineering team lead. Industries: manufacturing, automotive, mechanical engineering, industrial goods, consulting.

FACTS — MIM (Industrial Management, M.Sc.):
- Degree: M.Sc. Industrial Management. If the user means "International Management", note this MIM is INDUSTRIAL Management, not the bachelor International Management.
- Full-time, 3 semesters incl. thesis. 90 ECTS. 24 places. Starts summer and winter. Deadlines: 15 January / 15 June. Admission by selection procedure. Language of instruction: German and English.
- ENGLISH TRACK: completely taught in English. NO German certificate is required for admission — only English (e.g. TOEFL/IELTS). German is an optional course for daily life, not an admission requirement. (The standard German-taught MIM additionally needs German C1, only for non-native German speakers.) Never tell an English-Track applicant they need German; never say there is no English track. Page: https://techpf.hs-pforzheim.de/master/industrial_management_english_track
- Audience: interdisciplinary, for engineers and technical business-administration applicants.
- Themes: Technology & Innovation Management, Leadership, Process & Data Management, Networked Systems & Artificial Intelligence — i.e. managing innovation/tech-heavy organisations, people/leadership, data/process improvement, AI/networked systems in industry. Features: technical/business add-on qualifications, applied projects, many electives, specialization, possible semester abroad.
- Careers: innovation manager, technology/digital-transformation manager, business development manager, data/process analyst, product owner, consultant, project manager (industrial firms), broad engineering+business management. Industries: industrial firms, tech, consulting, automotive, logistics.
- MEM vs MIM: MEM = engineering-management, production/product/operations roles for industrial-engineering grads. MIM = broader — innovation, technology, data, general management (and has the English Track).

APPLICATION:
- Applications go through the HS Pforzheim online portal. In any answer about applying, include the application link as a clickable markdown link: [Master application page](https://www.hs-pforzheim.de/en/studies/study_programs/application_master). Where useful also link [MEM program page](https://techpf.hs-pforzheim.de/en/master_programs/engineering_and_management/engineering_and_management), [MIM program page](https://techpf.hs-pforzheim.de/en/master/industrial_management), and [MIM English Track](https://techpf.hs-pforzheim.de/master/industrial_management_english_track). Always use [label](url) syntax; never show a bare URL. For language-requirement questions, do NOT add a link line; just define the language requirements directly.
- When asked how to apply, about admission requirements, or about deadlines, STATE the actual application deadlines IN the answer: 15 January for the summer semester and 15 June for the winter semester. NEVER say "check the requirements/deadlines on the official website" or "see the application page for deadlines" — give the specific dates directly, not a pointer to go find them. (A clickable link may still be added afterwards as an extra, but the dates must be in the answer.)
- Typical documents: CV, university entrance qualification, first-degree transcript, motivation letter, recommendation, professional/educational proof, references, and language certificates as applicable. (Proof of stays abroad is NOT a required document — do not list it.) The bot estimates fit; HS Pforzheim makes the official decision.
- There is NO separate application fee for applying through the HS Pforzheim online portal. NEVER state, invent, or list an application/processing fee or a specific amount. (International applicants who must apply via uni-assist pay uni-assist's own handling fee — tell them to check uni-assist's current rate rather than quoting a number.) Do not add a "pay the fee" step.

INTERNATIONAL APPLICANTS (non-German degree):
- Need a recognized bachelor's equivalent to a German degree in a relevant engineering / industrial-engineering / technical-business field. Foreign qualifications usually checked via uni-assist, which issues a VPD (often obtained before applying).
- Language: English (IELTS/TOEFL) for English-taught programs — the ONLY language proof for the MIM English Track (an international student can enter it with no German at all). A German certificate is needed only for German-taught programs/parts, only for non-native German speakers whose prior degree was not in German; level is C1 (DSH-2, TestDaF TDN 4x4, telc C1 Hochschule, or Goethe C2 — C2 is one option, not the minimum). Native German speakers submit no German certificate.
- Visa/residence (non-EU): admission letter, blocked-account financial proof (amount set yearly by German authorities), health insurance. Some countries (e.g. India, China, Vietnam) also need an APS certificate. Documents usually need certified copies + English/German translations.
- Give concrete requirements first; you may add one line that exact amounts/scores and whether uni-assist/APS applies vary by country and should be confirmed per applicant.

TUITION & FEES (answer directly, never deflect):
- Baden-Württemberg charges tuition only to non-EU/non-EEA international students: 1,500 EUR/semester. EU/EEA students and German-Abitur holders pay none. Some are exempt (refugees, permanent residents, Erasmus/dual-degree); second-degree students pay 650 EUR.
- Everyone also pays a semester contribution (admin + student services) of ~150-170 EUR/semester, confirmed at enrolment.
- So: EU/EEA ≈ 150-170 EUR/semester only; non-EU = 1,500 EUR + that contribution. Fees can change by law — you may add one line to confirm at enrolment, but give the numbers first.

CONTACTS (MEM & MIM):
- Program leader: Prof. Dr. Ansgar Kuehn — Room T1.5.23, 07231 28 6490, ansgar.kuehn@hs-pforzheim.de
- Coordinator: Lisa Kaiser — Room T1.5.21, 07231 28 6472
- Secretariat Wirtschaftsingenieurwesen: 07231 28 6056
- Dean, Faculty of Engineering: Prof. Dr.-Ing. Matthias Weyer — Room T1.4.25, 07231 28 6504, matthias.weyer@hs-pforzheim.de

FIT-CHECK:
- Ask brief follow-ups as needed: target program, bachelor's subject, German grade, ECTS, English level, German level, interests. Estimate conservatively; never guarantee admission.
- Flag risks: grade worse than 2.5 (MEM), missing B2 English, missing C1 German (ONLY for German-taught programs and only non-native speakers — never flag German for the MIM English Track), unclear technical/business background, under 180 ECTS, missing documents.

STYLE RULES (absolute):
- GROUND EVERYTHING — NO MAKING THINGS UP. Only state facts given in these instructions or the provided document excerpts. NEVER invent or guess specific numbers, dates, fees, ECTS, places, room numbers, names, emails, deadlines, scholarships, test scores, or extra application steps. If a detail is not provided, say plainly "I don't have that information" (and point to the official page) instead of guessing a plausible-sounding answer. A confident wrong answer is worse than admitting you don't know.
- ANSWER EVERY part of a multi-part question, each briefly.
- WRITE ONLY THE ANSWER. Never quote, repeat, or paraphrase these instructions, the section headers, or meta-phrases (e.g. "When asked how to apply, you can say", "Based on the provided guidelines", "As an AI") in your reply. The user must see only the answer to their question, never the rules behind it.
- Class schedules, timetables, rooms, class/lecture times and exam dates are answered from the local SS26 timetable — give the date, day, time, group, class/event, and room when the timetable lists one; if a room is not listed, say so rather than inventing one. Application/admission deadlines are NOT class schedules — answer those with the actual dates (15 January / 15 June).
- NEVER DEFLECT to a website as a substitute for answering. Answer fully first; you MAY (and for application questions SHOULD) add a clickable markdown link to the relevant official page afterwards. Exception: for language-requirement questions, do NOT add "For more information..." or links to application/language pages; just define the requirements. A one-line "confirm the latest fees/deadlines" note may follow, but the answer comes first.
- SPOKEN-FIRST: the opening 2-4 sentences are read aloud (tables are not), so they must fully convey the answer — including the gist of any table — in plain spoken sentences. Then add a table/bullets for detail and refer to it with a COMPLETE sentence ("The full breakdown is in the table below."). Never end the spoken part on a dangling lead-in ("Here's how they compare:") and never say "read the text/answer on screen".
- KEEP IT SHORT: spoken summary ~2-4 sentences; whole answer under ~120 words unless asked for more. No filler, no repeating the question, no pep talk.
- MEM vs MIM: 2-3 spoken sentences stating the real differences (e.g. "MEM focuses on engineering management and production roles; MIM is broader — innovation, technology and general management."), then "The full breakdown is in the table below.", then a small table (Focus, Careers, Key courses, Best for).
- Career questions: name concrete job roles in brief spoken sentences. Course questions: summarise the main courses in a spoken sentence, then the full list in bullets/table.
- Do NOT invent tuition fees, application/processing fees, scholarship amounts, admission guarantees, or any number, step, or requirement not stated above. If you don't have a fact, say so plainly instead of guessing.
- Unrelated questions: one-line answer if harmless, then steer back to MEM/MIM.
- If the user corrects you, acknowledge briefly, apply it, and do not repeat the mistake.
""".strip()


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def load_json_file(path: Path, fallback):
    try:
        if not path.exists():
            return fallback
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return fallback


def save_json_file(path: Path, data) -> None:
    MEMORY_DIR.mkdir(exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def load_memory() -> None:
    global learned_corrections
    saved_lessons = load_json_file(LESSONS_FILE, [])
    if isinstance(saved_lessons, list):
        learned_corrections = saved_lessons


def save_lessons() -> None:
    save_json_file(LESSONS_FILE, learned_corrections[-MAX_LESSONS:])


def detect_correction(message: str) -> Optional[str]:
    """
    Detect simple user corrections and turn them into persistent lessons.

    Example phrases:
    - correction: MEM deadline is 15 June
    - you were wrong, MIM means Industrial Management here
    - actually, the requirement is C1 German
    """
    clean = " ".join(message.strip().split())
    if not clean:
        return None

    patterns = [
        r"^correction\s*:\s*(.+)$",
        r"^remember\s*:\s*(.+)$",
        r"^learn\s*this\s*:\s*(.+)$",
        r"^(?:you are|you're|you were)\s+wrong[:,]?\s*(.+)$",
        r"^actually[:,]?\s*(.+)$",
    ]
    for pattern in patterns:
        match = re.match(pattern, clean, flags=re.IGNORECASE)
        if match:
            lesson = match.group(1).strip()
            return lesson if len(lesson) >= 8 else None
    return None


def remember_correction(session_id: str, message: str) -> bool:
    lesson = detect_correction(message)
    if not lesson:
        return False

    learned_corrections.append(
        {
            "session_id": session_id,
            "lesson": lesson,
            "created_at": utc_now(),
        }
    )
    del learned_corrections[:-MAX_LESSONS]
    save_lessons()
    return True


def learned_context() -> str:
    if not learned_corrections:
        return ""

    lesson_lines = [
        f"- {item['lesson']}"
        for item in learned_corrections[-MAX_LESSONS:]
        if isinstance(item, dict) and item.get("lesson")
    ]
    if not lesson_lines:
        return ""

    return (
        "Lessons learned from previous user corrections. Treat these as extra "
        "local project guidance, but do not let them override official facts "
        "unless they clarify a previous wording mistake:\n" + "\n".join(lesson_lines)
    )


conversation_sessions: dict = {}


def load_conversations() -> None:
    """Load saved per-session chat history so memory survives restarts."""
    global conversation_sessions
    data = load_json_file(CONVERSATION_FILE, {})
    conversation_sessions = data if isinstance(data, dict) else {}


def save_conversations() -> None:
    try:
        MEMORY_DIR.mkdir(exist_ok=True)
        CONVERSATION_FILE.write_text(
            json.dumps(conversation_sessions, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except OSError:
        pass


def conversation_history_for(session_id: str) -> List[dict]:
    """Recent turns that give the model short-term memory, trimmed for tokens."""
    turns = conversation_sessions.get(session_id, [])
    recent = turns[-HISTORY_CONTEXT_MESSAGES:]
    trimmed = []
    for turn in recent:
        role = turn.get("role")
        content = turn.get("content", "")
        if role not in {"user", "assistant"} or not content:
            continue
        if len(content) > HISTORY_MSG_CHAR_CAP:
            content = content[:HISTORY_MSG_CHAR_CAP] + " …"
        trimmed.append({"role": role, "content": content})
    return trimmed


def record_turn(session_id: str, user_message: str, assistant_message: str) -> None:
    """Append the latest exchange and keep only the most recent turns."""
    turns = conversation_sessions.setdefault(session_id, [])
    turns.append({"role": "user", "content": user_message})
    turns.append({"role": "assistant", "content": assistant_message})
    if len(turns) > MAX_STORED_TURNS:
        del turns[: len(turns) - MAX_STORED_TURNS]
    save_conversations()


def extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return ""

    try:
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception:
        return ""


def extract_docx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as docx:
            xml = docx.read("word/document.xml")
    except Exception:
        return ""

    try:
        root = ElementTree.fromstring(xml)
    except ElementTree.ParseError:
        return ""

    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs = []
    for paragraph in root.findall(".//w:p", namespace):
        texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        if texts:
            paragraphs.append("".join(texts))
    return "\n".join(paragraphs)


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_texts.append(line)
            if slide_texts:
                texts.append(f"[Slide {slide_num}] " + " ".join(slide_texts))
        return "\n".join(texts)
    except ImportError:
        return ""
    except Exception:
        return ""


def extract_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix in {".pptx", ".ppt"}:
        return extract_pptx_text(path)
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return ""
    if suffix in {".html", ".htm"}:
        text = re.sub(r"<script\b[^>]*>.*?</script>", " ", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<style\b[^>]*>.*?</style>", " ", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"<[^>]+>", " ", text)
    return text


def split_document_chunks(text: str, source: str, chunk_size: int = 1200) -> List[dict]:
    clean = re.sub(r"\s+", " ", text).strip()
    if not clean:
        return []

    chunks = []
    start = 0
    while start < len(clean):
        end = min(start + chunk_size, len(clean))
        if end < len(clean):
            sentence_end = max(clean.rfind(". ", start, end), clean.rfind("\n", start, end))
            if sentence_end > start + 350:
                end = sentence_end + 1
        chunk_text = clean[start:end].strip()
        if chunk_text:
            chunks.append({"source": source, "text": chunk_text})
        start = max(end - 120, end)
    return chunks


def load_knowledge_documents() -> None:
    global knowledge_chunks
    KNOWLEDGE_DIR.mkdir(exist_ok=True)
    allowed_suffixes = {".pdf", ".docx", ".txt", ".md", ".html", ".htm", ".pptx", ".ppt"}
    chunks: List[dict] = []
    for path in sorted(KNOWLEDGE_DIR.rglob("*")):
        if not path.is_file() or path.suffix.lower() not in allowed_suffixes:
            continue
        # Keep the class timetable OUT of the chat knowledge — schedules belong
        # in the calendar UI, not in chat answers. (The calendar endpoint reads
        # TIMETABLE_FILE directly, so this does not affect the calendar.)
        if path.name in {"README.md", "source_urls.txt", "Semesteruebersicht_SS26.md"}:
            continue
        text = extract_document_text(path)
        chunks.extend(split_document_chunks(text, path.relative_to(ROOT).as_posix()))
    knowledge_chunks = chunks


def search_knowledge(query: str, limit: int = 2) -> List[dict]:
    terms = {
        term
        for term in re.findall(r"[a-zA-ZÄÖÜäöüß0-9]{3,}", query.lower())
        if term not in {"what", "when", "where", "which", "with", "about", "that", "this", "from"}
    }
    if not terms or not knowledge_chunks:
        return []

    ranked = []
    for chunk in knowledge_chunks:
        haystack = chunk["text"].lower()
        source = chunk["source"].lower()
        score = sum(haystack.count(term) * 2 + source.count(term) for term in terms)
        if score:
            ranked.append((score, chunk))

    ranked.sort(key=lambda item: item[0], reverse=True)
    return [chunk for _score, chunk in ranked[:limit]]


def knowledge_context(question: str) -> str:
    # Only the single best excerpt — keeps the per-request token cost low so we
    # stay under Groq's per-minute limit (the system prompt already covers the
    # core facts; this is just supplementary grounding).
    matches = search_knowledge(question, limit=1)
    if not matches:
        return ""

    lines = ["Relevant excerpts from local project documents. Prefer these when answering precise factual questions:"]
    for index, chunk in enumerate(matches, 1):
        lines.append(f"[Source {index}: {chunk['source']}]\n{chunk['text']}")
    return "\n\n".join(lines)


def contact_answer() -> str:
    return (
        "Here are the contact details for MEM and MIM at HS Pforzheim:\n\n"
        "| Role | Person / office | Details |\n"
        "|---|---|---|\n"
        "| Program leader for MEM and MIM | Prof. Dr. Ansgar Kuehn | Room T1.5.23, phone 07231 28 6490, email ansgar.kuehn@hs-pforzheim.de |\n"
        "| Program coordinator for MEM and MIM | Lisa Kaiser | Room T1.5.21, phone 07231 28 6472 |\n"
        "| Secretariat Wirtschaftsingenieurwesen | Program office | Phone 07231 28 6056 |\n"
        "| Dean of Faculty of Engineering | Prof. Dr.-Ing. Matthias Weyer | Room T1.4.25, phone 07231 28 6504, email matthias.weyer@hs-pforzheim.de |"
    )


def is_contact_question(question: str) -> bool:
    clean = question.lower()
    wants_contact = any(
        term in clean
        for term in (
            "contact",
            "email",
            "e-mail",
            "mail",
            "phone",
            "telephone",
            "dean",
            "dekan",
            "ansprech",
            "responsible person",
            "person responsible",
            "responsible",
            "coordinator",
            "program leader",
            "programme leader",
            "secretariat",
            "secretary",
            "who should i contact",
            "whom should i contact",
            "reach out",
            "get in touch",
        )
    )
    program_context = any(
        term in clean
        for term in (
            "mem",
            "mim",
            "engineering and management",
            "industrial management",
            "master",
            "masters",
            "program",
            "programme",
            "course",
            "study",
        )
    )
    asks_general_program_contact = any(
        term in clean
        for term in (
            "responsible",
            "coordinator",
            "program leader",
            "programme leader",
            "secretariat",
            "secretary",
            "dean",
        )
    )
    return wants_contact and (program_context or asks_general_program_contact)


def has_multiple_questions(question: str) -> bool:
    """Return True when the user is clearly asking more than one question."""
    if question.count("?") > 1:
        return True
    multi_patterns = [
        r"\band\s+(?:also\s+)?(?:what|when|where|how|which|who|is|are|do|does|can|could|will|would)\b",
        r"\balso\s+(?:what|when|where|how|which|who|is|are|do|does|can|could|will|would)\b",
        r"\band\s+(?:also\s+)?(?:tell me|explain|describe|give me)\b",
        r"\balso\s+(?:tell me|explain|describe|give me)\b",
    ]
    for pattern in multi_patterns:
        if re.search(pattern, question, flags=re.IGNORECASE):
            return True
    return False


def normalize_course_text(text: str) -> str:
    text = text.lower()
    text = text.replace("&", " and ")
    text = re.sub(r"\([^)]*\)", " ", text)
    text = re.sub(r"[^a-z0-9äöüß]+", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def course_lookup_tokens(text: str) -> set[str]:
    stopwords = {
        "who", "what", "which", "prof", "professor", "teacher", "lecturer",
        "instructor", "lehrperson", "teaches", "teach", "teaching", "course",
        "class", "for", "the", "is", "are", "my", "of", "and", "please",
    }
    return {
        token
        for token in normalize_course_text(text).split()
        if len(token) > 2 and token not in stopwords
    }


def load_course_professors() -> List[dict]:
    professors = []
    if not COURSE_PROFESSORS_FILE.exists():
        return professors
    for line in COURSE_PROFESSORS_FILE.read_text(encoding="utf-8").splitlines():
        if not line.startswith("| ") or line.startswith("|---") or "Course" in line:
            continue
        parts = [part.strip() for part in line.strip("|").split("|")]
        if len(parts) != 2:
            continue
        course, instructor = parts
        if course and instructor:
            professors.append({"course": course, "instructor": instructor})
    return professors


def direct_course_professor_answer(question: str) -> Optional[str]:
    clean = question.lower()
    asks_prof = any(
        term in clean
        for term in (
            "prof", "professor", "teacher", "lecturer", "instructor",
            "lehrperson", "who teaches", "who is teaching", "who takes",
        )
    )
    if not asks_prof:
        return None

    question_tokens = course_lookup_tokens(question)
    if not question_tokens:
        return "Which course do you mean? Ask for example: 'Who is the professor for Future Mobility?'"

    matches = []
    for item in load_course_professors():
        course_tokens = course_lookup_tokens(item["course"])
        overlap = question_tokens & course_tokens
        if not overlap:
            continue
        score = len(overlap) / max(len(course_tokens), 1)
        matches.append((score, len(overlap), item))

    if not matches:
        return "I don't have the professor information for that course."

    matches.sort(key=lambda match: (match[0], match[1], len(match[2]["course"])), reverse=True)
    best_score, _overlap, best = matches[0]
    if best_score < 0.45:
        return "I don't have the professor information for that course."

    return f"{best['course']} is taught by {best['instructor']}."


def direct_language_requirements_answer(question: str) -> Optional[str]:
    clean = question.lower()
    # Only fire on clear REQUIREMENT questions. Do NOT trigger on bare profile
    # tokens (ielts, c1, b2, ...) — those also appear when a user states their own
    # background in a fit-check ("I have IELTS 6.5, C1 German"), which must reach
    # the fit-check / LLM path, not this canned language answer.
    language_terms = (
        "language requirement", "language requirements",
        "english requirement", "german requirement",
        "german certificate", "english certificate",
        "german required", "english required",
        "is german required", "is english required",
        "do i need german", "do i need english",
        "what language", "which language",
    )
    if not any(term in clean for term in language_terms):
        return None

    mentions_mem = "mem" in clean or "engineering and management" in clean
    mentions_mim = "mim" in clean or "industrial management" in clean
    mentions_english_track = any(
        term in clean
        for term in ("english track", "mim english", "mim e-track", "mim e track", "mim3")
    )

    if mentions_english_track:
        return (
            "For the MIM English Track, only English proof is required, such as "
            "TOEFL or IELTS. No German certificate is required for admission. "
            "German is only offered as an optional course for daily life."
        )

    if mentions_mem:
        return (
            "For MEM, you need English at B2 level. German C1 is required only if "
            "you are not a native German speaker and your previous degree was not "
            "taught in German. Accepted German examples are DSH-2, TestDaF TDN "
            "4x4, or telc C1 Hochschule; C2 is not the minimum."
        )

    if mentions_mim:
        return (
            "For the MIM English Track, only English proof is required, such as "
            "TOEFL or IELTS; no German certificate is required. For the standard "
            "German-taught MIM, German C1 is required only for non-native German "
            "speakers whose previous degree was not taught in German."
        )

    return (
        "Language requirements depend on the program. MEM requires English B2, "
        "plus German C1 only for non-native German speakers whose previous degree "
        "was not taught in German. MIM English Track requires only English proof "
        "such as TOEFL or IELTS; no German certificate is required. Standard "
        "German-taught MIM requires German C1 where applicable."
    )


def direct_precise_answer(question: str) -> Optional[str]:
    if has_multiple_questions(question):
        return None

    professor_answer = direct_course_professor_answer(question)
    if professor_answer:
        return professor_answer

    language_answer = direct_language_requirements_answer(question)
    if language_answer:
        return language_answer

    fit_answer = direct_fit_check_answer(question)
    if fit_answer:
        return fit_answer

    timetable_answer = direct_timetable_answer(question)
    if timetable_answer:
        return timetable_answer

    clean = question.lower()
    asks_track = any(term in clean for term in ("english track", "english-taught", "english taught", "taught in english"))
    mentions_mim = "mim" in clean or "industrial management" in clean
    # If the question is really about German/language/eligibility detail, skip
    # this short blurb and let the LLM answer with the full grounded facts.
    asks_language_detail = any(
        term in clean
        for term in (
            "german", "deutsch", "certificate", "certification", "c1", "c2", "b2",
            "requirement", "require", "ielts", "toefl", "language", "level", "eligib",
        )
    )
    if asks_track and mentions_mim and not asks_language_detail:
        return (
            "Yes — MIM has an English Track that is taught completely in English.\n\n"
            "International students can join it with English proof only (e.g. TOEFL or "
            "IELTS); no German certificate is required for admission."
        )

    if is_contact_question(question):
        return contact_answer()
    return None


def direct_fit_check_answer(question: str) -> Optional[str]:
    clean = question.lower()
    fit_terms = (
        "course match",
        "program match",
        "programme match",
        "fit check",
        "eligibility",
        "eligible",
        "am i suitable",
        "am i a good fit",
        "which program fits",
        "which programme fits",
        "which course fits",
    )
    if not any(term in clean for term in fit_terms):
        return None

    return (
        "Sure. I can check whether MEM or MIM looks like the better match for you.\n\n"
        "Please send these details:\n\n"
        "1. Your bachelor's degree subject\n"
        "2. Your final grade or current average\n"
        "3. Your ECTS credits\n"
        "4. Your English level\n"
        "5. Your German level\n"
        "6. Your main interest: engineering management, production/product strategy, technology innovation, data/process management, AI/networked systems, or leadership\n\n"
        "Quick guide: MEM usually fits applicants with an industrial engineering / Wirtschaftsingenieurwesen background and strong interest in engineering-management roles. MIM is broader for engineers and technical business applicants, including the English-taught track."
    )


_DEFLECTION_PHRASES = (
    "visit the official",
    "check the official",
    "checking the official",
    "consult the official",
    "refer to the official",
    "navigate to the",
    "visit the hs pforzheim",
    "visit the university",
    "visit the website",
    "visit their website",
    "please visit",
    "recommend visiting",
    "i recommend you visit",
    "try checking",
    "i suggest checking",
    "please check the official",
    "you should check the official",
    "i don't have the specific contact",
    "i do not have the specific contact",
    "for the most up-to-date information, please",
    "for the most current information, please",
    "for accurate and current information",
    "official website for",
    "website for the latest",
    "website for more",
)


def _strip_deflection_sentences(answer: str) -> str:
    parts = re.split(r"(?<=[.!?])\s+(?=[A-Z\-\*])", answer)
    kept = [p for p in parts if not any(t in p.lower() for t in _DEFLECTION_PHRASES)]
    result = " ".join(kept).strip()
    return result if len(result) > 60 else answer


def replace_unhelpful_deflection(question: str, answer: str) -> str:
    lower_answer = answer.lower()
    deflects = any(phrase in lower_answer for phrase in _DEFLECTION_PHRASES)

    if is_contact_question(question) and deflects:
        return contact_answer()

    if deflects:
        return _strip_deflection_sentences(answer)

    return answer


def parse_timetable_rows() -> List[dict]:
    files = []
    if TIMETABLE_FILE.exists():
        files.append(TIMETABLE_FILE)
    files.extend(path for path in sorted(KNOWLEDGE_DIR.glob(TIMETABLE_FILE_PATTERN)) if path != TIMETABLE_FILE)

    rows = []
    for path in files:
        for line in path.read_text(encoding="utf-8").splitlines():
            if not line.startswith("| 20"):
                continue
            parts = [part.strip() for part in line.strip("|").split("|")]
            if len(parts) not in {6, 7}:
                continue
            row = {
                "date": parts[0],
                "day": parts[1],
                "block": parts[2],
                "time": parts[3],
                "audience": parts[4],
                "entry": parts[5],
                "room": "",
                "source": path.name,
            }
            if len(parts) == 7:
                row["room"] = parts[6]
            rows.append(row)

    # Some classes appear in both the overview and a room-aware personal plan.
    # Keep one copy, preferring the row that has a room.
    deduped = {}
    for row in rows:
        normalized_entry = re.sub(r"\s*\([^)]*\)\s*$", "", row["entry"]).strip().lower()
        key = (row["date"], row["time"], row["audience"], normalized_entry)
        existing = deduped.get(key)
        if existing is None or (row.get("room") and not existing.get("room")):
            deduped[key] = row
    return sorted(deduped.values(), key=block_sort_key)


def extract_question_date(question: str) -> Optional[str]:
    clean = question.strip()
    lowered = clean.lower()
    today = datetime.now().date()
    if re.search(r"\btoday\b|\bheute\b", lowered):
        return today.strftime("%Y-%m-%d")
    if re.search(r"\btomorrow\b|\bmorgen\b", lowered):
        return (today + timedelta(days=1)).strftime("%Y-%m-%d")

    iso_match = re.search(r"\b(20\d{2})-(\d{1,2})-(\d{1,2})\b", clean)
    if iso_match:
        year, month, day = iso_match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"

    month_names = {
        "january": 1,
        "jan": 1,
        "february": 2,
        "feb": 2,
        "march": 3,
        "mar": 3,
        "märz": 3,
        "maerz": 3,
        "april": 4,
        "apr": 4,
        "may": 5,
        "mai": 5,
        "june": 6,
        "jun": 6,
        "juni": 6,
        "july": 7,
        "jul": 7,
        "juli": 7,
        "august": 8,
        "aug": 8,
        "september": 9,
        "sep": 9,
        "october": 10,
        "oct": 10,
        "oktober": 10,
        "november": 11,
        "nov": 11,
        "december": 12,
        "dec": 12,
        "dezember": 12,
    }
    pattern = r"\b(\d{1,2})(?:st|nd|rd|th|\.)?\s+([A-Za-zÄÖÜäöüß]+)\s+(20\d{2})\b"
    match = re.search(pattern, clean, flags=re.IGNORECASE)
    if not match:
        return None
    day, month_name, year = match.groups()
    month = month_names.get(month_name.lower())
    if not month:
        return None
    return f"{int(year):04d}-{month:02d}-{int(day):02d}"


def extract_day_of_month(question: str) -> Optional[int]:
    match = re.search(r"\b(\d{1,2})(?:st|nd|rd|th|\.)?\b", question.lower())
    if not match:
        return None
    day = int(match.group(1))
    if 1 <= day <= 31:
        return day
    return None


WEEKDAY_NAMES = {
    "monday": "Montag",
    "mondays": "Montag",
    "montag": "Montag",
    "dienstag": "Dienstag",
    "tuesday": "Dienstag",
    "tuesdays": "Dienstag",
    "wednesday": "Mittwoch",
    "wednesdays": "Mittwoch",
    "mittwoch": "Mittwoch",
    "thursday": "Donnerstag",
    "thursdays": "Donnerstag",
    "donnerstag": "Donnerstag",
    "friday": "Freitag",
    "fridays": "Freitag",
    "freitag": "Freitag",
    "saturday": "Samstag",
    "saturdays": "Samstag",
    "samstag": "Samstag",
    "sunday": "Sonntag",
    "sundays": "Sonntag",
    "sonntag": "Sonntag",
}


def extract_question_weekday(question: str) -> Optional[str]:
    clean = question.lower()
    for word, german_day in WEEKDAY_NAMES.items():
        if re.search(rf"\b{re.escape(word)}(?:'s)?\b", clean):
            return german_day
    return None


def extract_audience(question: str) -> Optional[str]:
    clean = question.lower()
    if (
        "mim e-track" in clean
        or "mim e track" in clean
        or "english track" in clean
        or "mim english" in clean
        or "mim3" in clean
        or "mim 3" in clean
        or "varsha" in clean
        or "anuka" in clean
    ):
        return "MIM E-Track"
    if "mim / 1" in clean or "mim 1" in clean or "mim first" in clean:
        return "MIM / 1"
    if "mim / 2" in clean or "mim 2" in clean or "mim second" in clean:
        return "MIM / 2"
    if "mem / 1" in clean or "mem 1" in clean or "mem first" in clean:
        return "MEM / 1"
    if "mem / 2" in clean or "mem 2" in clean or "mem second" in clean:
        return "MEM / 2"
    return None


def infer_personal_audience(question: str) -> Optional[str]:
    clean = question.lower()
    if re.search(r"\bmy\b.{0,40}\b(schedule|timetable|class|classes|lecture|lectures)\b", clean):
        return "MIM E-Track"
    personal_terms = (
        "my schedule", "my timetable", "my class", "my classes", "my lecture",
        "my lectures", "what do i have", "where should i go", "where do i go",
        "which room should i go", "for me",
    )
    if any(term in clean for term in personal_terms):
        return "MIM E-Track"
    return None


def row_matches_audience(row: dict, audience: Optional[str]) -> bool:
    if not audience:
        return True
    row_audience = row.get("audience", "")
    return row_audience == audience or audience in [part.strip() for part in row_audience.split(",")]


def block_sort_key(row: dict) -> tuple[str, int, str]:
    block_text = str(row.get("block", ""))
    match = re.search(r"\d+", block_text)
    block = int(match.group(0)) if match else 999
    return str(row.get("date", "")), block, str(row.get("time", ""))


def nearest_timetable_date_for_weekday(rows: List[dict], weekday: str) -> Optional[str]:
    today = datetime.now().date()
    future_dates = sorted(
        {
            row["date"]
            for row in rows
            if row.get("day") == weekday
            and datetime.strptime(row["date"], "%Y-%m-%d").date() >= today
        }
    )
    if future_dates:
        return future_dates[0]
    all_dates = sorted({row["date"] for row in rows if row.get("day") == weekday})
    return all_dates[-1] if all_dates else None


def resolve_timetable_date(question: str, rows: List[dict], weekday: Optional[str]) -> Optional[str]:
    explicit_date = extract_question_date(question)
    if explicit_date:
        return explicit_date

    day = extract_day_of_month(question)
    if day:
        candidates = []
        today = datetime.now().date()
        for row in rows:
            try:
                row_date = datetime.strptime(row["date"], "%Y-%m-%d").date()
            except (KeyError, ValueError):
                continue
            if row_date.day != day:
                continue
            if weekday and row.get("day") != weekday:
                continue
            candidates.append(row_date)

        if candidates:
            future = sorted(date for date in candidates if date >= today)
            return (future[0] if future else sorted(candidates)[-1]).strftime("%Y-%m-%d")

    if weekday:
        return nearest_timetable_date_for_weekday(rows, weekday)
    return None


def course_tokens(question: str) -> List[str]:
    stopwords = {
        "which", "what", "where", "when", "room", "number", "should", "go",
        "class", "lecture", "schedule", "timetable", "today", "tomorrow",
        "monday", "tuesday", "wednesday", "thursday", "friday", "saturday",
        "sunday", "montag", "dienstag", "mittwoch", "donnerstag", "freitag",
        "samstag", "sonntag", "mem", "mim", "track", "english", "first",
        "second", "for", "the", "and", "with", "have", "has", "please",
    }
    return [
        token
        for token in re.findall(r"[a-zA-ZÄÖÜäöüß][a-zA-ZÄÖÜäöüß-]{2,}", question.lower())
        if token not in stopwords
    ]


def filter_rows_by_course(rows: List[dict], question: str) -> List[dict]:
    tokens = course_tokens(question)
    if not tokens:
        return rows

    scored = []
    for row in rows:
        haystack = row.get("entry", "").lower()
        score = sum(1 for token in tokens if token in haystack)
        if score:
            scored.append((score, row))

    if not scored:
        return []
    best_score = max(score for score, _row in scored)
    return [row for score, row in scored if score == best_score]


def entry_tokens(entry: str) -> set[str]:
    stopwords = {
        "and", "und", "der", "die", "das", "for", "mit", "von", "bis",
        "kurs", "course", "prof", "dr", "lb", "ab", "uhr",
    }
    return {
        token
        for token in re.findall(r"[a-zA-ZÄÖÜäöüß][a-zA-ZÄÖÜäöüß-]{2,}", entry.lower())
        if token not in stopwords
    }


def entry_similarity(left: str, right: str) -> int:
    return len(entry_tokens(left) & entry_tokens(right))


def same_timetable_subject(left: str, right: str) -> bool:
    left_tokens = entry_tokens(left)
    right_tokens = entry_tokens(right)
    return len(left_tokens & right_tokens) >= 2 or ("capstone" in left_tokens and "capstone" in right_tokens)


def prefer_roomed_rows(rows: List[dict]) -> List[dict]:
    roomed_rows = [row for row in rows if row.get("room")]
    if not roomed_rows:
        return rows

    filtered = []
    for row in rows:
        if row.get("room"):
            filtered.append(row)
            continue
        has_roomed_match = any(
            other.get("date") == row.get("date")
            and other.get("audience") == row.get("audience")
            and same_timetable_subject(other.get("entry", ""), row.get("entry", ""))
            for other in roomed_rows
        )
        if not has_roomed_match:
            filtered.append(row)
    return filtered


def format_timetable_rows(rows: List[dict], title: str, audience: Optional[str]) -> str:
    rows = sorted(prefer_roomed_rows(rows), key=block_sort_key)
    lines = [title]
    for row in rows[:14]:
        room = row.get("room") or "not listed in the timetable"
        lines.append(
            "- Time: {time} | {entry} [{audience}] | Date: {date} ({day}) | Block: {block} | Room: {room}".format(
                date=row.get("date", ""),
                day=row.get("day", ""),
                time=row.get("time") or "time not listed",
                block=row.get("block") or "-",
                entry=row.get("entry", ""),
                audience=row.get("audience", ""),
                room=room,
            )
        )
    if len(rows) > 14:
        lines.append(f"- Plus {len(rows) - 14} more matching entries. Add your group or course name to narrow it down.")
    if not audience:
        lines.append("Tip: include your group, for example MEM / 1 or MIM E-Track, to get only your own classes.")
    return "\n".join(lines)


def direct_timetable_answer(question: str) -> Optional[str]:
    """Answer schedule, weekday, class-time, exam-date, and room questions from
    the local SS26 timetable."""
    clean = question.lower()

    # Never hijack admission/application questions (e.g. "application deadline").
    if any(term in clean for term in ("application", "admission", "apply", "applying", "enrol", "enroll")):
        return None

    phrase_terms = (
        "timetable", "time table", "schedule", "my class", "my classes",
        "my lecture", "my exam", "what do i have", "what's on", "whats on",
        "what is on", "which room", "what room", "room number", "class time",
        "where should i go", "where do i go", "this week", "next week",
        "tomorrow", "today", "tonight",
    )
    has_phrase = any(term in clean for term in phrase_terms)
    has_weekday = extract_question_weekday(question) is not None
    has_exam = re.search(r"\bexams?\b", clean) is not None
    has_date = extract_question_date(question) is not None

    if not (has_phrase or has_weekday or has_exam or has_date):
        return None

    rows = parse_timetable_rows()
    if not rows:
        return "I cannot find the local SS26 timetable file right now."

    audience = extract_audience(question) or infer_personal_audience(question)
    rows = [row for row in rows if row_matches_audience(row, audience)]

    weekday = extract_question_weekday(question)
    target_date = resolve_timetable_date(question, rows, weekday)
    if target_date:
        day_rows = [row for row in rows if row.get("date") == target_date]
        if day_rows:
            course_rows = filter_rows_by_course(day_rows, question)
            if course_rows:
                day_rows = course_rows
            label = f"Your schedule for {target_date}" if infer_personal_audience(question) else f"Schedule for {target_date}"
            return format_timetable_rows(day_rows, f"{label}:", audience)
        group_note = f" for {audience}" if audience else ""
        return f"I do not see any timetable entries{group_note} on {target_date}."

    if weekday:
        target_date = resolve_timetable_date(question, rows, weekday)
        if not target_date:
            return f"I do not see any {weekday} entries in the SS26 timetable."
        day_rows = [row for row in rows if row.get("date") == target_date]
        return format_timetable_rows(day_rows, f"{weekday} schedule for {target_date}:", audience)

    course_rows = filter_rows_by_course(rows, question)
    if course_rows and course_rows != rows:
        today = datetime.now().date()
        future_rows = [
            row for row in course_rows
            if datetime.strptime(row["date"], "%Y-%m-%d").date() >= today
        ]
        selected = future_rows or course_rows
        return format_timetable_rows(selected[:14], "Matching timetable entries:", audience)

    return (
        "I can answer that from the SS26 timetable, but I need a day/date or class name. "
        "For example: 'Monday schedule for MEM / 1' or 'Which room for Future Mobility?'"
    )


load_memory()
load_conversations()
load_knowledge_documents()


class ChatRequest(BaseModel):
    message: str = Field(..., min_length=1, max_length=2000)
    session_id: str = Field(default="default", max_length=120)
    # "explore" = prospective student deciding; "enrolled" = current student.
    mode: str = Field(default="explore", max_length=20)


class ChatResponse(BaseModel):
    text: str
    audio_base64: str
    audio_mime: str = AUDIO_MIME
    voice_id: str
    llm_model: str
    learned: bool = False


class VoiceRequest(BaseModel):
    text: str = Field(..., min_length=1, max_length=4000)


class VoiceResponse(BaseModel):
    audio_base64: str
    audio_mime: str = AUDIO_MIME
    voice_id: str


app = FastAPI(title="MEM & MIM Guide Bot API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)


@app.get("/", response_class=HTMLResponse)
async def index() -> str:
    """Serve the single-page chatbot frontend."""
    if not HTML_FILE.exists():
        raise HTTPException(status_code=404, detail="HTML file not found.")
    return HTML_FILE.read_text(encoding="utf-8")


@app.get("/LOGO.png")
async def logo() -> FileResponse:
    """Serve the Hochschule Pforzheim logo used in the page header."""
    logo_file = ROOT / "LOGO.png"
    if not logo_file.exists():
        raise HTTPException(status_code=404, detail="Logo file not found.")
    return FileResponse(logo_file, media_type="image/png")


@app.get("/health")
async def health() -> dict:
    """Small health check endpoint for testing the server."""
    return {
        "ok": True,
        "provider": "groq" if GROQ_API_KEY else "openrouter",
        "model": GROQ_MODEL if GROQ_API_KEY else OPENROUTER_MODEL,
        "elevenlabs_voice_id": ELEVENLABS_VOICE_ID,
        "voice_provider": VOICE_PROVIDER,
        "xtts_reference_clip": str(resolve_project_path(XTTS_SPEAKER_WAV)),
        "knowledge_documents": len({chunk["source"] for chunk in knowledge_chunks}),
        "knowledge_chunks": len(knowledge_chunks),
        "learned_corrections": len(learned_corrections),
    }


@app.get("/api/lessons")
async def lessons() -> dict:
    """Return corrections the bot has learned from."""
    return {"lessons": learned_corrections[-MAX_LESSONS:]}


@app.post("/api/reload-knowledge")
async def reload_knowledge() -> dict:
    """Reload local PDF/DOCX/TXT/MD/HTML files from knowledge_docs."""
    load_knowledge_documents()
    return {
        "ok": True,
        "knowledge_documents": len({chunk["source"] for chunk in knowledge_chunks}),
        "knowledge_chunks": len(knowledge_chunks),
    }


@app.get("/api/timetable")
async def timetable() -> dict:
    """Return the SS26 timetable rows + group list for the calendar popover."""
    rows = parse_timetable_rows()
    groups = sorted({row["audience"] for row in rows})
    return {"rows": rows, "groups": groups}


@app.post("/api/chat", response_model=ChatResponse)
async def chat(payload: ChatRequest) -> ChatResponse:
    """Generate a grounded AI answer quickly. Voice is generated on demand."""
    question = payload.message.strip()
    if not question:
        raise HTTPException(status_code=400, detail="Please send a message.")

    learned = remember_correction(payload.session_id, question)
    history = conversation_history_for(payload.session_id)
    answer = direct_precise_answer(question)
    if not answer:
        if not GROQ_API_KEY and not OPENROUTER_API_KEY:
            raise HTTPException(status_code=500, detail="Set GROQ_API_KEY or OPENROUTER_API_KEY in your .env file.")
        answer = await ask_ai(question, payload.mode, history)
    answer = replace_unhelpful_deflection(question, answer)
    if learned:
        answer = "Got it. I saved that correction and will use it in future answers.\n\n" + answer
    # Remember this exchange so the next message can build on it.
    record_turn(payload.session_id, question, answer)

    active_model = GROQ_MODEL if GROQ_API_KEY else OPENROUTER_MODEL
    return ChatResponse(
        text=answer,
        audio_base64="",
        audio_mime=AUDIO_MIME,
        voice_id=active_voice_id(),
        llm_model=active_model,
        learned=learned,
    )


@app.post("/api/voice", response_model=VoiceResponse)
async def voice(payload: VoiceRequest) -> VoiceResponse:
    """Generate professor-style audio only when the user clicks the speaker."""
    audio_b64, audio_mime = await synthesize_voice_audio(payload.text)
    return VoiceResponse(
        audio_base64=audio_b64,
        audio_mime=audio_mime,
        voice_id=active_voice_id(),
    )


def _file_to_b64(path: Path) -> tuple[str, str]:
    return base64.b64encode(path.read_bytes()).decode("ascii"), audio_mime_for(path)


async def synthesize_voice_audio(answer: str) -> tuple[str, str]:
    """
    Generate professor-style audio for THIS specific answer.

      cartesia      ->  Cartesia hosted clone (fast, ~1s)  -> local XTTS
      openvoice_hf  ->  HuggingFace OpenVoice V2            -> local XTTS
      xtts          ->  local XTTS clone

    If real per-answer synthesis fails we return EMPTY audio on purpose. The
    browser then reads the actual answer text with its own voice. We do NOT fall
    back to the static professor sample clip here, because that would make every
    answer play the same recording — which sounds like the bot is broken.
    """
    if VOICE_PROVIDER in {"browser", "none", "off"}:
        return "", ""

    if VOICE_PROVIDER in {"sample", "demo"}:
        sample = find_professor_demo_audio_file()
        if sample:
            return _file_to_b64(sample)
        return "", ""

    # Primary: hosted cloning (fast, runs on their GPUs). No lock needed —
    # stateless network calls, so prefetches can overlap.
    if VOICE_PROVIDER == "cartesia":
        try:
            audio_file = await synthesize_cartesia(answer)
            return _file_to_b64(audio_file)
        except Exception as exc:
            print(f"Cartesia failed, falling back to local XTTS: {exc}")

    if VOICE_PROVIDER == "fish":
        try:
            audio_file = await synthesize_fish(answer)
            return _file_to_b64(audio_file)
        except Exception as exc:
            print(f"Fish failed, falling back to local XTTS: {exc}")

    if VOICE_PROVIDER in {"elevenlabs", "11labs"}:
        try:
            audio_file = await synthesize_elevenlabs_file(answer)
            return _file_to_b64(audio_file)
        except Exception as exc:
            print(f"ElevenLabs failed, falling back to local XTTS: {exc}")

    # Primary: HuggingFace OpenVoice V2 (clones professor voice on HF servers).
    if VOICE_PROVIDER in {"openvoice_hf", "openvoice"}:
        try:
            audio_file = await synthesize_openvoice_hf(answer)
            return _file_to_b64(audio_file)
        except Exception as exc:
            print(f"OpenVoice HF failed, falling back to local XTTS: {exc}")

    # Local XTTS clone (works fully offline). Serialized: the shared model is not
    # concurrency-safe, so only one local synthesis runs at a time.
    if VOICE_PROVIDER in {"cartesia", "fish", "elevenlabs", "11labs", "openvoice_hf", "openvoice", "xtts"}:
        try:
            async with voice_lock:
                audio_file = await synthesize_xtts(answer)
            return _file_to_b64(audio_file)
        except Exception as exc:
            print(f"XTTS synthesis failed for this answer; browser voice will read it: {exc}")

    # Deliberately return empty so the browser reads the real answer text,
    # instead of replaying an unrelated static clip for every answer.
    return "", ""


def active_voice_id() -> str:
    if VOICE_PROVIDER == "cartesia":
        return f"cartesia:{CARTESIA_MODEL}"
    if VOICE_PROVIDER == "fish":
        return "fish-audio-professor-clone"
    if VOICE_PROVIDER in {"elevenlabs", "11labs"}:
        return f"elevenlabs:{ELEVENLABS_VOICE_ID}"
    if VOICE_PROVIDER in {"openvoice_hf", "openvoice"}:
        return f"hf-openvoice:{OPENVOICE_HF_SPACE}"
    if VOICE_PROVIDER == "xtts":
        return "local-xtts-professor-reference"
    if VOICE_PROVIDER in {"sample", "demo"}:
        return "local-professor-demo-sample"
    if VOICE_PROVIDER in {"browser", "none", "off"}:
        return "browser-speech"
    return ELEVENLABS_VOICE_ID


def audio_mime_for(path: Path) -> str:
    if path.suffix.lower() == ".wav":
        return "audio/wav"
    if path.suffix.lower() == ".m4a":
        return "audio/mp4"
    return AUDIO_MIME


def find_professor_demo_audio_file() -> Optional[Path]:
    if PROFESSOR_DEMO_AUDIO_FILE:
        configured = Path(PROFESSOR_DEMO_AUDIO_FILE)
        return configured if configured.exists() else None

    for filename in ("professor_voice_sample.mp3", "professor_voice_sample.wav", "professor_voice_sample.m4a"):
        candidate = ROOT / "voice_source/generated" / filename
        if candidate.exists():
            return candidate
    return None


def resolve_project_path(path_value: str) -> Path:
    path = Path(path_value)
    return path if path.is_absolute() else ROOT / path


def clean_text_for_speech(text: str) -> str:
    """
    Turn a chat answer into natural spoken text.

    The bot answers conversationally first, then may add a table. Tables do not
    read well aloud, so we speak only the conversational prose (which already
    says things like "see the table below"). We never read out table rows or
    append "please read the screen".
    """
    spoken_lines = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        # Skip markdown table rows and divider rows.
        if line.startswith("|") or re.match(r"^\|?\s*:?-{3,}", line):
            continue
        if line.count("|") >= 2:
            continue
        spoken_lines.append(line)

    clean = " ".join(spoken_lines) if spoken_lines else text
    clean = re.sub(r"`([^`]+)`", r"\1", clean)
    clean = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", clean)
    clean = clean.replace("|", " ").replace("#", "")
    clean = re.sub(r"[*_~>]+", "", clean)
    clean = re.sub(r"\s+", " ", clean).strip()

    # Fast hosted voices (ElevenLabs/Cartesia/Fish) can read the whole answer;
    # only the slow local XTTS needs a tight cap to stay quick.
    limit = HOSTED_MAX_CHARS if VOICE_PROVIDER in {"elevenlabs", "11labs", "cartesia", "fish"} else XTTS_MAX_CHARS
    if len(clean) <= limit:
        return clean
    # Cut at a sentence end so the spoken text ends naturally, not mid-word.
    clipped = clean[:limit]
    sentence_end = max(clipped.rfind(". "), clipped.rfind("! "), clipped.rfind("? "))
    if sentence_end > 80:
        return clipped[: sentence_end + 1].strip()
    return clipped.rsplit(" ", 1)[0].strip()


async def fish_create_model() -> str:
    """Clone the professor clip into a Fish Audio model and return its id."""
    speaker = resolve_project_path(XTTS_SPEAKER_WAV)
    if not speaker.exists():
        raise RuntimeError(f"Reference clip not found: {speaker}")

    headers = {"Authorization": f"Bearer {FISH_API_KEY}"}
    async with httpx.AsyncClient(timeout=180) as client:
        files = {"voices": (speaker.name, speaker.read_bytes(), "audio/mpeg")}
        data = {"title": "MEM MIM Professor", "type": "tts", "train_mode": "fast", "visibility": "private"}
        response = await client.post("https://api.fish.audio/model", headers=headers, files=files, data=data)
    if response.status_code >= 400:
        raise RuntimeError(f"Fish clone failed: {response.status_code} {response.text}")
    model_id = response.json().get("_id")
    if not model_id:
        raise RuntimeError(f"Fish clone returned no model id: {response.text}")
    return model_id


async def get_fish_reference_id() -> str:
    if FISH_REFERENCE_ID:
        return FISH_REFERENCE_ID
    if FISH_REFERENCE_ID_FILE.exists():
        cached = FISH_REFERENCE_ID_FILE.read_text(encoding="utf-8").strip()
        if cached:
            return cached
    model_id = await fish_create_model()
    FISH_REFERENCE_ID_FILE.parent.mkdir(parents=True, exist_ok=True)
    FISH_REFERENCE_ID_FILE.write_text(model_id, encoding="utf-8")
    print(f"Cloned professor voice on Fish Audio: {model_id}")
    return model_id


async def fish_tts_bytes(text: str, reference_id: str) -> bytes:
    headers = {"Authorization": f"Bearer {FISH_API_KEY}", "Content-Type": "application/json"}
    body = {"text": text, "reference_id": reference_id, "format": FISH_FORMAT}
    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.post("https://api.fish.audio/v1/tts", headers=headers, json=body)
    if response.status_code >= 400:
        raise RuntimeError(f"Fish TTS failed: {response.status_code} {response.text}")
    if not response.content:
        raise RuntimeError("Fish returned empty audio.")
    return response.content


async def synthesize_fish(text: str) -> Path:
    """Clone the professor voice via Fish Audio's hosted API. Cached by text."""
    if not FISH_API_KEY:
        raise RuntimeError("FISH_API_KEY is not set.")
    speech_text = clean_text_for_speech(text)
    if not speech_text:
        raise RuntimeError("No text available for Fish.")

    reference_id = await get_fish_reference_id()
    cache_key = hashlib.sha256(
        f"fish|{reference_id}|{FISH_FORMAT}|{speech_text}".encode("utf-8")
    ).hexdigest()[:20]
    suffix = ".mp3" if FISH_FORMAT == "mp3" else ".wav"
    output = OPENVOICE_OUTPUT_DIR / f"fish_{cache_key}{suffix}"
    if output.exists():
        return output

    audio = await fish_tts_bytes(speech_text, reference_id)
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary_output = output.with_suffix(".tmp" + suffix)
    temporary_output.write_bytes(audio)
    temporary_output.replace(output)
    return output


async def synthesize_elevenlabs_file(text: str) -> Path:
    """ElevenLabs TTS in the professor voice, cached by text to save credits."""
    if not ELEVENLABS_API_KEY:
        raise RuntimeError("ELEVENLABS_API_KEY is not set.")
    speech_text = clean_text_for_speech(text)
    if not speech_text:
        raise RuntimeError("No text available for ElevenLabs.")

    cache_key = hashlib.sha256(
        f"elevenlabs|{ELEVENLABS_VOICE_ID}|{ELEVENLABS_MODEL_ID}|{ELEVENLABS_LANGUAGE}|{speech_text}".encode("utf-8")
    ).hexdigest()[:20]
    output = OPENVOICE_OUTPUT_DIR / f"el_{cache_key}.mp3"
    if output.exists():
        return output

    audio = await synthesize_elevenlabs(speech_text)
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary_output = output.with_suffix(".tmp.mp3")
    temporary_output.write_bytes(audio)
    temporary_output.replace(output)
    return output


async def cartesia_clone_voice() -> str:
    """Clone the professor reference clip on Cartesia and return a voice id."""
    speaker = resolve_project_path(XTTS_SPEAKER_WAV)
    if not speaker.exists():
        raise RuntimeError(f"Reference clip not found: {speaker}")

    headers = {"X-API-Key": CARTESIA_API_KEY, "Cartesia-Version": CARTESIA_VERSION}
    async with httpx.AsyncClient(timeout=180) as client:
        files = {"clip": (speaker.name, speaker.read_bytes(), "audio/mpeg")}
        data = {
            "name": "MEM MIM Professor",
            "description": "Professor reference voice for the MEM/MIM guide bot.",
            "language": CARTESIA_LANGUAGE,
            "mode": "similarity",
            "enhance": "true",
        }
        response = await client.post(
            "https://api.cartesia.ai/voices/clone", headers=headers, files=files, data=data
        )
    if response.status_code >= 400:
        raise RuntimeError(f"Cartesia clone failed: {response.status_code} {response.text}")
    voice_id = response.json().get("id")
    if not voice_id:
        raise RuntimeError(f"Cartesia clone returned no voice id: {response.text}")
    return voice_id


async def get_cartesia_voice_id() -> str:
    """Return the configured/cached Cartesia voice id, cloning once if needed."""
    if CARTESIA_VOICE_ID:
        return CARTESIA_VOICE_ID
    if CARTESIA_VOICE_ID_FILE.exists():
        cached = CARTESIA_VOICE_ID_FILE.read_text(encoding="utf-8").strip()
        if cached:
            return cached
    voice_id = await cartesia_clone_voice()
    CARTESIA_VOICE_ID_FILE.parent.mkdir(parents=True, exist_ok=True)
    CARTESIA_VOICE_ID_FILE.write_text(voice_id, encoding="utf-8")
    print(f"Cloned professor voice on Cartesia: {voice_id}")
    return voice_id


async def cartesia_tts_bytes(text: str, voice_id: str) -> bytes:
    headers = {
        "X-API-Key": CARTESIA_API_KEY,
        "Cartesia-Version": CARTESIA_VERSION,
        "Content-Type": "application/json",
    }
    body = {
        "model_id": CARTESIA_MODEL,
        "transcript": text,
        "voice": {"mode": "id", "id": voice_id},
        "output_format": {"container": "wav", "encoding": "pcm_s16le", "sample_rate": 44100},
        "language": CARTESIA_LANGUAGE,
    }
    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.post("https://api.cartesia.ai/tts/bytes", headers=headers, json=body)
    if response.status_code >= 400:
        raise RuntimeError(f"Cartesia TTS failed: {response.status_code} {response.text}")
    if not response.content:
        raise RuntimeError("Cartesia returned empty audio.")
    return response.content


async def synthesize_cartesia(text: str) -> Path:
    """Clone the professor voice fast via Cartesia's hosted API. Cached by text."""
    if not CARTESIA_API_KEY:
        raise RuntimeError("CARTESIA_API_KEY is not set.")
    speech_text = clean_text_for_speech(text)
    if not speech_text:
        raise RuntimeError("No text available for Cartesia.")

    voice_id = await get_cartesia_voice_id()
    cache_key = hashlib.sha256(
        f"cartesia|{CARTESIA_MODEL}|{voice_id}|{speech_text}".encode("utf-8")
    ).hexdigest()[:20]
    output = OPENVOICE_OUTPUT_DIR / f"ca_{cache_key}.wav"
    if output.exists():
        return output

    audio = await cartesia_tts_bytes(speech_text, voice_id)
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary_output = output.with_suffix(".tmp.wav")
    temporary_output.write_bytes(audio)
    temporary_output.replace(output)
    return output


async def synthesize_openvoice_hf(text: str) -> Path:
    """Clone the professor voice via the HuggingFace OpenVoice V2 Space."""
    speech_text = clean_text_for_speech(text)
    if not speech_text:
        raise RuntimeError("No text available for OpenVoice.")

    speaker = resolve_project_path(XTTS_SPEAKER_WAV)
    if not speaker.exists():
        raise RuntimeError(f"Reference clip not found: {speaker}")

    cache_key = hashlib.sha256(
        f"openvoice|{OPENVOICE_HF_SPACE}|{OPENVOICE_STYLE}|{speaker}|{speech_text}".encode("utf-8")
    ).hexdigest()[:20]
    output = OPENVOICE_OUTPUT_DIR / f"ov_{cache_key}.wav"
    if output.exists():
        return output

    return await asyncio.to_thread(run_openvoice_hf_to_file, speech_text, speaker, output)


def run_openvoice_hf_to_file(text: str, speaker: Path, output: Path) -> Path:
    global openvoice_hf_client

    output.parent.mkdir(parents=True, exist_ok=True)

    from gradio_client import Client

    # gradio_client >=1.0 needs handle_file() for file inputs; older versions
    # accept a plain filepath string. Support both.
    try:
        from gradio_client import handle_file
        reference = handle_file(str(speaker))
    except ImportError:
        reference = str(speaker)

    if openvoice_hf_client is None:
        openvoice_hf_client = Client(OPENVOICE_HF_SPACE)

    result = openvoice_hf_client.predict(
        text,              # Text Prompt
        OPENVOICE_STYLE,   # Style, e.g. en_default
        reference,         # Reference Audio (professor clip)
        True,              # Agree checkbox
        fn_index=OPENVOICE_FN_INDEX,
    )

    # Returns (info, synthesised_audio_path, reference_audio_path).
    audio_path = result[1] if isinstance(result, (list, tuple)) and len(result) > 1 else result
    if not audio_path or not Path(audio_path).exists():
        raise RuntimeError(f"OpenVoice Space returned no usable audio: {result!r}")

    import shutil
    shutil.copyfile(audio_path, output)
    return output


async def synthesize_xtts(text: str) -> Path:
    speech_text = clean_text_for_speech(text)
    if not speech_text:
        raise RuntimeError("No text available for XTTS.")

    speaker = resolve_project_path(XTTS_SPEAKER_WAV)
    if not speaker.exists():
        raise RuntimeError(f"XTTS reference clip not found: {speaker}")

    cache_key = hashlib.sha256(
        f"{XTTS_LANGUAGE}|{speaker}|{speech_text}".encode("utf-8")
    ).hexdigest()[:20]
    output = XTTS_OUTPUT_DIR / f"{cache_key}.wav"
    if output.exists():
        return output

    return await asyncio.to_thread(run_xtts_to_file, speech_text, speaker, output)


def run_xtts_to_file(text: str, speaker: Path, output: Path) -> Path:
    global xtts_model

    output.parent.mkdir(parents=True, exist_ok=True)
    os.environ.setdefault("COQUI_TOS_AGREED", "1")
    os.environ.setdefault("TTS_HOME", str(ROOT / "models"))

    if xtts_model is None:
        from TTS.api import TTS

        xtts_model = TTS("tts_models/multilingual/multi-dataset/xtts_v2", gpu=False)

    temporary_output = output.with_suffix(".tmp.wav")
    xtts_model.tts_to_file(
        text=text,
        speaker_wav=str(speaker),
        language=XTTS_LANGUAGE,
        file_path=str(temporary_output),
    )
    temporary_output.replace(output)
    return output


def configured_ai_providers() -> List[dict]:
    providers = []
    if GROQ_API_KEY:
        providers.append(
            {
                "name": "Groq",
                "url": "https://api.groq.com/openai/v1/chat/completions",
                "model": GROQ_MODEL,
                "headers": {
                    "Authorization": f"Bearer {GROQ_API_KEY}",
                    "Content-Type": "application/json",
                },
            }
        )
    if GEMINI_API_KEY:
        providers.append(
            {
                "name": "Gemini",
                "url": "https://generativelanguage.googleapis.com/v1beta/openai/chat/completions",
                "model": GEMINI_MODEL,
                "headers": {
                    "Authorization": f"Bearer {GEMINI_API_KEY}",
                    "Content-Type": "application/json",
                },
            }
        )
    if OPENROUTER_API_KEY and OPENROUTER_ENABLED:
        providers.append(
            {
                "name": "OpenRouter",
                "url": "https://openrouter.ai/api/v1/chat/completions",
                "model": OPENROUTER_MODEL,
                "headers": {
                    "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                    "Content-Type": "application/json",
                    "HTTP-Referer": os.getenv("HTTP_REFERER", "http://localhost:8000"),
                    "X-Title": "MEM MIM Guide Bot",
                },
            }
        )
    return providers


def mode_context(mode: str) -> str:
    if (mode or "").lower() == "enrolled":
        return (
            "The user is a CURRENTLY ENROLLED MEM/MIM student. Focus on student "
            "life: timetable, which classes/lectures they have, rooms, times, "
            "exam dates, and deadlines, using the SS26 timetable. If you don't "
            "know their group, ask once which one they are: MEM/1, MEM/2, MIM/1, "
            "MIM/2, or MIM E-Track. Keep answers short and practical."
        )
    return (
        "The user is a PROSPECTIVE student exploring MEM/MIM and deciding if it "
        "fits them. Focus on guiding their decision: program differences, "
        "courses, careers, requirements, application steps, and fit based on "
        "their background and interests."
    )


async def ask_ai(question: str, mode: str = "explore", history: Optional[List[dict]] = None) -> str:
    """Try Groq first when configured, then fall back to OpenRouter."""
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "system", "content": mode_context(mode)},
    ]
    local_knowledge = knowledge_context(question)
    if local_knowledge:
        messages.append({"role": "system", "content": local_knowledge})
    local_lessons = learned_context()
    if local_lessons:
        messages.append({"role": "system", "content": local_lessons})
    # Prior turns give the bot short-term memory so it can handle follow-ups.
    if history:
        messages.extend(history)
    messages.append({"role": "user", "content": question})

    providers = configured_ai_providers()
    last_error = ""

    async with httpx.AsyncClient(timeout=60) as client:
        for index, provider in enumerate(providers):
            is_last_provider = index == len(providers) - 1
            request_body = {
                "model": provider["model"],
                "messages": messages,
                "temperature": 0.1,
                "max_tokens": 400,
            }

            # Retry transparently on per-minute rate limits (free Groq tier has a
            # low TPM limit; a quick wait + retry hides it from the user).
            for attempt in range(4):
                response = await client.post(
                    provider["url"],
                    headers=provider["headers"],
                    json=request_body,
                )

                if response.status_code == 429:
                    # If another provider is available (e.g. Gemini), hand off to
                    # it immediately instead of waiting — that keeps answers fast.
                    if not is_last_provider:
                        last_error = f"{provider['name']} rate limited; switching to fallback"
                        print(last_error)
                        break
                    match = re.search(r"try again in ([\d.]+)\s*s", response.text)
                    wait_s = min(float(match.group(1)) + 0.5, 8.0) if match else 4.0
                    last_error = f"{provider['name']} rate limited; retrying in {wait_s:.1f}s"
                    print(last_error)
                    await asyncio.sleep(wait_s)
                    continue

                if response.status_code >= 400:
                    last_error = f"{provider['name']} request failed: {response.text}"
                    print(last_error)
                    break  # non-rate-limit error -> try next provider

                data = response.json()
                try:
                    text = data["choices"][0]["message"]["content"].strip()
                except (KeyError, IndexError, AttributeError):
                    last_error = f"{provider['name']} returned no text."
                    print(last_error)
                    break

                if text:
                    return text
                last_error = f"{provider['name']} returned an empty response."
                print(last_error)
                break

    # All providers failed (usually Groq's per-minute rate limit). Surface a calm,
    # user-friendly message instead of the raw API error. last_error is logged.
    print(f"All AI providers failed: {last_error}")
    raise HTTPException(
        status_code=503,
        detail="I'm getting a lot of questions right now and hit a brief limit. Please wait a few seconds and ask again.",
    )


async def synthesize_elevenlabs(text: str) -> bytes:
    """
    Convert the model's text answer to speech with ElevenLabs.

    The voice is controlled by ELEVENLABS_VOICE_ID. By default it uses the
    standard pre-made placeholder voice ID 21m00Tcm4TlvDq8ikWAM, which you can
    replace with the professor's cloned voice ID after cloning is complete.
    """
    url = (
        f"https://api.elevenlabs.io/v1/text-to-speech/{ELEVENLABS_VOICE_ID}"
        f"?output_format={ELEVENLABS_OUTPUT_FORMAT}"
    )
    request_body = {
        "text": text,
        "model_id": ELEVENLABS_MODEL_ID,
        "voice_settings": {
            "stability": 0.45,
            "similarity_boost": 0.75,
            "style": 0.25,
            "use_speaker_boost": True,
        },
    }
    # Force English so the voice can't drift into German. Embedded German words
    # (names, "Hochschule Pforzheim", etc.) are still spoken, just inside an
    # English base. Only sent for models that actually honor the lock.
    if ELEVENLABS_LANGUAGE and ELEVENLABS_MODEL_ID in ELEVENLABS_LANGUAGE_LOCK_MODELS:
        request_body["language_code"] = ELEVENLABS_LANGUAGE

    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.post(
            url,
            headers={
                "xi-api-key": ELEVENLABS_API_KEY,
                "Content-Type": "application/json",
                "Accept": AUDIO_MIME,
            },
            json=request_body,
        )

    if response.status_code >= 400:
        raise HTTPException(
            status_code=502,
            detail=f"ElevenLabs request failed: {response.text}",
        )

    if not response.content:
        raise HTTPException(status_code=502, detail="ElevenLabs returned empty audio.")
    return response.content


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(
        "app:app",
        host="127.0.0.1",
        port=int(os.getenv("PORT", "8000")),
        reload=True,
    )
